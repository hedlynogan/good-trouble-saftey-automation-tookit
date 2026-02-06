#!/usr/bin/env python3
"""
Generate a PowerPoint guide: How to Create an Emergency iCloud Contact List
with Apple Contacts for the Good Trouble Safety Automation Toolkit.
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# ── Palette ──────────────────────────────────────────────────────────────
DARK_BG       = RGBColor(0x1A, 0x1A, 0x2E)
ACCENT_BLUE   = RGBColor(0x00, 0x7A, 0xFF)
ACCENT_RED    = RGBColor(0xE0, 0x3E, 0x3E)
WHITE         = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY    = RGBColor(0xCC, 0xCC, 0xCC)
MED_GRAY      = RGBColor(0x99, 0x99, 0x99)
DARK_TEXT      = RGBColor(0x33, 0x33, 0x33)
SOFT_BG       = RGBColor(0xF5, 0xF7, 0xFA)
STEP_BG       = RGBColor(0xE8, 0xF0, 0xFE)
TIP_BG        = RGBColor(0xFF, 0xF3, 0xCD)
WARN_BG       = RGBColor(0xFD, 0xE8, 0xE8)
GREEN         = RGBColor(0x28, 0xA7, 0x45)


def set_slide_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_textbox(slide, left, top, width, height, text,
                font_size=18, bold=False, color=WHITE,
                alignment=PP_ALIGN.LEFT, font_name="Calibri"):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font_name
    p.alignment = alignment
    return txBox


def add_bullet_list(slide, left, top, width, height, items,
                    font_size=16, color=WHITE, spacing=Pt(6),
                    font_name="Calibri", bold_first=False):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = item
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.name = font_name
        p.space_after = spacing
        if bold_first and i == 0:
            p.font.bold = True
        p.level = 0
    return txBox


def add_rounded_rect(slide, left, top, width, height, fill_color):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    return shape


def add_step_card(slide, number, title, instructions, top_offset,
                  bg_color=STEP_BG):
    """Add a numbered step card with a colored badge and instructions."""
    card_left = Inches(0.6)
    card_width = Inches(8.8)
    card_height = Inches(1.1)

    # Card background
    add_rounded_rect(slide, card_left, top_offset, card_width, card_height, bg_color)

    # Number badge
    badge = slide.shapes.add_shape(
        MSO_SHAPE.OVAL, card_left + Inches(0.15),
        top_offset + Inches(0.15), Inches(0.55), Inches(0.55)
    )
    badge.fill.solid()
    badge.fill.fore_color.rgb = ACCENT_BLUE
    badge.line.fill.background()
    tf = badge.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.text = str(number)
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.font.name = "Calibri"
    p.alignment = PP_ALIGN.CENTER
    tf.paragraphs[0].space_before = Pt(0)
    tf.paragraphs[0].space_after = Pt(0)

    # Title
    add_textbox(slide, card_left + Inches(0.85), top_offset + Inches(0.08),
                Inches(7.5), Inches(0.35), title,
                font_size=17, bold=True, color=DARK_TEXT)

    # Instructions
    add_textbox(slide, card_left + Inches(0.85), top_offset + Inches(0.42),
                Inches(7.7), Inches(0.65), instructions,
                font_size=13, color=DARK_TEXT)


def add_tip_box(slide, text, top, icon_label="TIP", bg=TIP_BG, label_color=RGBColor(0x85, 0x6D, 0x0)):
    left = Inches(0.6)
    width = Inches(8.8)
    height = Inches(0.7)
    add_rounded_rect(slide, left, top, width, height, bg)
    add_textbox(slide, left + Inches(0.2), top + Inches(0.05),
                Inches(0.7), Inches(0.3), icon_label,
                font_size=11, bold=True, color=label_color)
    add_textbox(slide, left + Inches(0.2), top + Inches(0.28),
                Inches(8.3), Inches(0.4), text,
                font_size=12, color=DARK_TEXT)


# ═════════════════════════════════════════════════════════════════════════
# BUILD PRESENTATION
# ═════════════════════════════════════════════════════════════════════════
prs = Presentation()
prs.slide_width = Inches(10)
prs.slide_height = Inches(7.5)
blank_layout = prs.slide_layouts[6]  # blank

# ─── SLIDE 1 – Title ─────────────────────────────────────────────────────
slide = prs.slides.add_slide(blank_layout)
set_slide_bg(slide, DARK_BG)

# Top accent line
add_rounded_rect(slide, Inches(0), Inches(0), Inches(10), Inches(0.06), ACCENT_BLUE)

add_textbox(slide, Inches(0.8), Inches(1.5), Inches(8.4), Inches(1.0),
            "How to Create an Emergency\niCloud Contact List",
            font_size=36, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER,
            font_name="Calibri")

add_textbox(slide, Inches(1.0), Inches(3.0), Inches(8.0), Inches(0.6),
            "Using Apple Contacts on iPhone, iPad, or Mac",
            font_size=22, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)

# Divider
add_rounded_rect(slide, Inches(3.5), Inches(3.9), Inches(3.0), Inches(0.04), ACCENT_BLUE)

add_textbox(slide, Inches(1.0), Inches(4.3), Inches(8.0), Inches(0.5),
            "Good Trouble Safety Automation Toolkit",
            font_size=18, color=ACCENT_BLUE, alignment=PP_ALIGN.CENTER,
            font_name="Calibri")

add_textbox(slide, Inches(1.0), Inches(5.0), Inches(8.0), Inches(0.8),
            "Set up a dedicated contact list so your emergency alerts\n"
            "reach the right people instantly.",
            font_size=15, color=MED_GRAY, alignment=PP_ALIGN.CENTER)

add_textbox(slide, Inches(0.5), Inches(6.8), Inches(9.0), Inches(0.4),
            "Supports iOS 15+ | macOS Ventura+ | iCloud Sync",
            font_size=12, color=MED_GRAY, alignment=PP_ALIGN.CENTER)

# ─── SLIDE 2 – Overview / Why ────────────────────────────────────────────
slide = prs.slides.add_slide(blank_layout)
set_slide_bg(slide, SOFT_BG)

add_textbox(slide, Inches(0.6), Inches(0.4), Inches(8.8), Inches(0.6),
            "Why Create an Emergency Contact List?",
            font_size=28, bold=True, color=DARK_TEXT, alignment=PP_ALIGN.LEFT)

add_rounded_rect(slide, Inches(0.6), Inches(0.95), Inches(3.0), Inches(0.04), ACCENT_BLUE)

reasons = [
    "Quickly send alerts to trusted people during an encounter",
    "Keep emergency contacts organized and separate from your full address book",
    "Syncs automatically across all your Apple devices via iCloud",
    "Works seamlessly with iOS Shortcuts for the Good Trouble toolkit",
    "Update the list once and changes propagate everywhere",
]
for i, reason in enumerate(reasons):
    y = Inches(1.4) + Inches(i * 0.7)
    # Bullet icon rect
    add_rounded_rect(slide, Inches(0.8), y + Inches(0.05), Inches(0.25), Inches(0.25), ACCENT_BLUE)
    add_textbox(slide, Inches(0.82), y + Inches(0.02), Inches(0.25), Inches(0.3),
                ">>", font_size=10, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(1.2), y, Inches(8.0), Inches(0.5),
                reason, font_size=16, color=DARK_TEXT)

add_tip_box(slide,
            "The contact list you create here is what the Good Trouble Shortcut uses "
            "to know who to message when you trigger an emergency alert.",
            Inches(5.5))

# ─── SLIDE 3 – Prerequisites ─────────────────────────────────────────────
slide = prs.slides.add_slide(blank_layout)
set_slide_bg(slide, SOFT_BG)

add_textbox(slide, Inches(0.6), Inches(0.4), Inches(8.8), Inches(0.6),
            "Before You Begin",
            font_size=28, bold=True, color=DARK_TEXT)
add_rounded_rect(slide, Inches(0.6), Inches(0.95), Inches(2.0), Inches(0.04), ACCENT_BLUE)

prereqs = [
    ("iCloud Account", "Sign in to iCloud on your device (Settings > [Your Name] > iCloud)."),
    ("Contacts Syncing ON", "Under iCloud settings, make sure Contacts toggle is enabled."),
    ("Apple Contacts App", "Pre-installed on every iPhone, iPad, and Mac."),
    ("Emergency Contact Info", "Gather phone numbers / emails of your trusted contacts."),
]
for i, (title, desc) in enumerate(prereqs):
    y = Inches(1.3) + Inches(i * 0.95)
    # Number circle
    badge = slide.shapes.add_shape(
        MSO_SHAPE.OVAL, Inches(0.8), y + Inches(0.05),
        Inches(0.45), Inches(0.45)
    )
    badge.fill.solid()
    badge.fill.fore_color.rgb = ACCENT_BLUE
    badge.line.fill.background()
    tf = badge.text_frame
    p = tf.paragraphs[0]
    p.text = str(i + 1)
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER

    add_textbox(slide, Inches(1.45), y, Inches(7.5), Inches(0.35),
                title, font_size=17, bold=True, color=DARK_TEXT)
    add_textbox(slide, Inches(1.45), y + Inches(0.35), Inches(7.5), Inches(0.4),
                desc, font_size=13, color=DARK_TEXT)

add_tip_box(slide,
            "If iCloud Contacts syncing is off, your list will only exist on that one device "
            "and won't be available to Shortcuts on your other devices.",
            Inches(5.4), icon_label="NOTE", bg=WARN_BG, label_color=ACCENT_RED)

# ─── SLIDE 4 – Step-by-step: Create the List (iPhone/iPad) ───────────────
slide = prs.slides.add_slide(blank_layout)
set_slide_bg(slide, SOFT_BG)

add_textbox(slide, Inches(0.6), Inches(0.3), Inches(8.8), Inches(0.6),
            "Create the List on iPhone or iPad",
            font_size=28, bold=True, color=DARK_TEXT)
add_rounded_rect(slide, Inches(0.6), Inches(0.85), Inches(3.5), Inches(0.04), ACCENT_BLUE)

add_step_card(slide, 1, "Open the Contacts App",
              "Tap the Contacts app on your Home Screen (or open Phone and tap the Contacts tab).",
              Inches(1.1))

add_step_card(slide, 2, "Go to the Lists View",
              "Tap \"Lists\" in the upper-left corner to see all your contact lists. "
              "If you're already on the Lists screen, skip this step.",
              Inches(2.35))

add_step_card(slide, 3, "Create a New List",
              "Tap \"Add List\" in the upper-left. Name it  Emergency Contacts  (or any name you prefer). "
              "Make sure iCloud is selected as the account, then tap \"Done\".",
              Inches(3.6))

add_step_card(slide, 4, "Add Contacts to the List",
              "Tap your new list to open it. Tap the \"+\" button or \"Add Contact\" to browse "
              "your existing contacts. Select each trusted person, then tap \"Done\".",
              Inches(4.85))

add_tip_box(slide,
            "You can also drag and drop contacts into the list on a Mac using the Contacts app sidebar.",
            Inches(6.2))

# ─── SLIDE 5 – Step-by-step: Create the List (Mac) ───────────────────────
slide = prs.slides.add_slide(blank_layout)
set_slide_bg(slide, SOFT_BG)

add_textbox(slide, Inches(0.6), Inches(0.3), Inches(8.8), Inches(0.6),
            "Create the List on Mac",
            font_size=28, bold=True, color=DARK_TEXT)
add_rounded_rect(slide, Inches(0.6), Inches(0.85), Inches(2.5), Inches(0.04), ACCENT_BLUE)

add_step_card(slide, 1, "Open Contacts",
              "Launch the Contacts app from Applications, the Dock, or Spotlight (Cmd + Space, type \"Contacts\").",
              Inches(1.1))

add_step_card(slide, 2, "Show the Sidebar",
              "If you don't see the sidebar with lists, go to View > Show Groups (or click the sidebar icon).",
              Inches(2.35))

add_step_card(slide, 3, "Create a New List",
              "Click File > New List (or press Cmd + Shift + N in older macOS). "
              "Name it  Emergency Contacts  and make sure it appears under the iCloud section.",
              Inches(3.6))

add_step_card(slide, 4, "Add Contacts",
              "Select contacts from \"All Contacts\", then drag them onto the new list in the sidebar. "
              "You can hold Cmd to select multiple contacts at once.",
              Inches(4.85))

add_tip_box(slide,
            "On macOS Ventura and later, the Contacts app has a refreshed sidebar. "
            "Right-click in the sidebar to see \"New List\" as well.",
            Inches(6.2))

# ─── SLIDE 6 – Verify iCloud Sync ────────────────────────────────────────
slide = prs.slides.add_slide(blank_layout)
set_slide_bg(slide, SOFT_BG)

add_textbox(slide, Inches(0.6), Inches(0.3), Inches(8.8), Inches(0.6),
            "Verify iCloud Sync",
            font_size=28, bold=True, color=DARK_TEXT)
add_rounded_rect(slide, Inches(0.6), Inches(0.85), Inches(2.5), Inches(0.04), ACCENT_BLUE)

add_step_card(slide, 1, "Check on Another Device",
              "Open Contacts on a second Apple device signed into the same iCloud account. "
              "Your \"Emergency Contacts\" list should appear automatically.",
              Inches(1.1))

add_step_card(slide, 2, "Check on iCloud.com",
              "Go to iCloud.com in a browser, sign in, and open Contacts. "
              "Your list should appear in the sidebar under \"Lists\".",
              Inches(2.35))

add_step_card(slide, 3, "Force a Sync (if needed)",
              "On iPhone/iPad: Settings > [Your Name] > iCloud > toggle Contacts off and back on. "
              "On Mac: System Settings > Apple ID > iCloud > toggle Contacts.",
              Inches(3.6))

add_tip_box(slide,
            "iCloud syncing may not be instant. Give it a minute or two, "
            "and make sure you're connected to Wi-Fi or cellular data.",
            Inches(5.0))

add_tip_box(slide,
            "If the list shows up under \"On My iPhone\" instead of \"iCloud\", you need to move "
            "the contacts. Delete the local list and recreate it while iCloud is the default account "
            "(Settings > Contacts > Default Account > iCloud).",
            Inches(5.9), icon_label="WARN", bg=WARN_BG, label_color=ACCENT_RED)

# ─── SLIDE 7 – Using with Good Trouble Shortcut ──────────────────────────
slide = prs.slides.add_slide(blank_layout)
set_slide_bg(slide, SOFT_BG)

add_textbox(slide, Inches(0.6), Inches(0.3), Inches(8.8), Inches(0.6),
            "Using the List with Good Trouble Shortcuts",
            font_size=28, bold=True, color=DARK_TEXT)
add_rounded_rect(slide, Inches(0.6), Inches(0.85), Inches(4.0), Inches(0.04), ACCENT_BLUE)

add_textbox(slide, Inches(0.6), Inches(1.1), Inches(8.8), Inches(0.6),
            "Your emergency contact list integrates directly with the iOS Shortcut workflow:",
            font_size=15, color=DARK_TEXT)

add_step_card(slide, 1, "Open the Good Trouble Shortcut in the Shortcuts App",
              "Find the shortcut you built following the ShortcutTutorial. Tap the \"...\" to edit it.",
              Inches(1.6))

add_step_card(slide, 2, "Locate the \"Send Message\" or \"Get Contacts\" Action",
              "In the shortcut, find the action that retrieves contacts to message. "
              "It may say \"Get Contacts from [list]\" or \"Send Message to [contacts]\".",
              Inches(2.85))

add_step_card(slide, 3, "Point It to Your Emergency Contacts List",
              "Tap the list name in the action and choose your \"Emergency Contacts\" list. "
              "The shortcut will now message everyone on that list when triggered.",
              Inches(4.1))

add_step_card(slide, 4, "Test the Shortcut",
              "Run the shortcut manually to verify it picks up all the contacts from your list. "
              "Confirm that each person receives the alert message with your GPS location.",
              Inches(5.35))

add_tip_box(slide,
            "When you add or remove someone from the Emergency Contacts list, the shortcut "
            "automatically picks up the change - no need to edit the shortcut again.",
            Inches(6.55))

# ─── SLIDE 8 – Best Practices ────────────────────────────────────────────
slide = prs.slides.add_slide(blank_layout)
set_slide_bg(slide, SOFT_BG)

add_textbox(slide, Inches(0.6), Inches(0.3), Inches(8.8), Inches(0.6),
            "Best Practices",
            font_size=28, bold=True, color=DARK_TEXT)
add_rounded_rect(slide, Inches(0.6), Inches(0.85), Inches(1.8), Inches(0.04), ACCENT_BLUE)

practices = [
    ("Keep the List Small and Trusted",
     "3-5 people is ideal. These should be people who will act on an alert immediately."),
    ("Include Varied Contact Methods",
     "Ensure contacts have both phone numbers and email addresses saved in their card."),
    ("Inform Your Contacts",
     "Let everyone on the list know they're your emergency contacts and what an alert looks like."),
    ("Review Regularly",
     "Check the list every few months. Remove outdated contacts, add new trusted people."),
    ("Test Periodically",
     "Run a test alert to make sure messages go through and contacts know how to respond."),
]
for i, (title, desc) in enumerate(practices):
    y = Inches(1.1) + Inches(i * 0.85)
    # Green check circle
    badge = slide.shapes.add_shape(
        MSO_SHAPE.OVAL, Inches(0.8), y + Inches(0.05),
        Inches(0.4), Inches(0.4)
    )
    badge.fill.solid()
    badge.fill.fore_color.rgb = GREEN
    badge.line.fill.background()
    tf = badge.text_frame
    p = tf.paragraphs[0]
    p.text = "+"
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER

    add_textbox(slide, Inches(1.4), y, Inches(7.8), Inches(0.3),
                title, font_size=16, bold=True, color=DARK_TEXT)
    add_textbox(slide, Inches(1.4), y + Inches(0.32), Inches(7.8), Inches(0.4),
                desc, font_size=13, color=DARK_TEXT)

# ─── SLIDE 9 – Troubleshooting ───────────────────────────────────────────
slide = prs.slides.add_slide(blank_layout)
set_slide_bg(slide, SOFT_BG)

add_textbox(slide, Inches(0.6), Inches(0.3), Inches(8.8), Inches(0.6),
            "Troubleshooting",
            font_size=28, bold=True, color=DARK_TEXT)
add_rounded_rect(slide, Inches(0.6), Inches(0.85), Inches(2.2), Inches(0.04), ACCENT_RED)

issues = [
    ("List doesn't appear on other devices",
     "Verify iCloud Contacts sync is ON on all devices. Check that you're signed in with the same Apple ID. "
     "Go to Settings > [Your Name] > iCloud > Show All > Contacts (toggle ON)."),
    ("List shows under \"On My iPhone\" instead of iCloud",
     "Set your default account to iCloud: Settings > Contacts > Default Account > iCloud. "
     "Then recreate the list so it's stored in iCloud."),
    ("Shortcut says \"No contacts found\"",
     "Make sure the shortcut action references the correct list name (exact spelling). "
     "Re-select the list in the shortcut's \"Get Contacts\" action."),
    ("Contacts missing from the list",
     "Open the list, tap \"Add Contact\", and re-add them. Contacts must be explicitly added "
     "to the list - they don't appear automatically."),
]
for i, (problem, solution) in enumerate(issues):
    y = Inches(1.15) + Inches(i * 1.25)
    add_rounded_rect(slide, Inches(0.6), y, Inches(8.8), Inches(1.1), RGBColor(0xFF, 0xFF, 0xFF))
    # Red indicator bar
    add_rounded_rect(slide, Inches(0.6), y, Inches(0.08), Inches(1.1), ACCENT_RED)
    add_textbox(slide, Inches(0.9), y + Inches(0.08), Inches(8.3), Inches(0.35),
                problem, font_size=15, bold=True, color=ACCENT_RED)
    add_textbox(slide, Inches(0.9), y + Inches(0.45), Inches(8.3), Inches(0.6),
                solution, font_size=12, color=DARK_TEXT)

# ─── SLIDE 10 – Summary / Closing ────────────────────────────────────────
slide = prs.slides.add_slide(blank_layout)
set_slide_bg(slide, DARK_BG)

add_rounded_rect(slide, Inches(0), Inches(0), Inches(10), Inches(0.06), ACCENT_BLUE)

add_textbox(slide, Inches(0.8), Inches(1.2), Inches(8.4), Inches(0.8),
            "You're All Set!",
            font_size=34, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)

add_rounded_rect(slide, Inches(3.5), Inches(2.1), Inches(3.0), Inches(0.04), ACCENT_BLUE)

summary_items = [
    "Your emergency iCloud contact list is created and syncing",
    "Your Good Trouble Shortcut is connected to the list",
    "Adding or removing contacts updates the alert automatically",
    "Test your setup periodically to make sure it works",
]
for i, item in enumerate(summary_items):
    y = Inches(2.5) + Inches(i * 0.55)
    add_textbox(slide, Inches(1.5), y, Inches(7.0), Inches(0.4),
                f"  {item}", font_size=16, color=LIGHT_GRAY,
                alignment=PP_ALIGN.CENTER)

add_textbox(slide, Inches(0.8), Inches(5.0), Inches(8.4), Inches(0.8),
            "For full setup instructions, see the ShortcutTutorial\n"
            "and WorkflowTutorial guides in the doc/ios/ folder.",
            font_size=14, color=MED_GRAY, alignment=PP_ALIGN.CENTER)

add_textbox(slide, Inches(0.8), Inches(6.2), Inches(8.4), Inches(0.5),
            "Good Trouble Safety Automation Toolkit",
            font_size=16, bold=True, color=ACCENT_BLUE, alignment=PP_ALIGN.CENTER)

add_textbox(slide, Inches(0.8), Inches(6.7), Inches(8.4), Inches(0.4),
            "Stay safe. Stay informed. Know your rights.",
            font_size=13, color=MED_GRAY, alignment=PP_ALIGN.CENTER)

# ═════════════════════════════════════════════════════════════════════════
# SAVE
# ═════════════════════════════════════════════════════════════════════════
output_path = (
    "/Volumes/OWC STX NVMe/Develop/Source/GoodTrouble/"
    "good-trouble-saftey-automation-tookit/doc/ios/"
    "Emergency_iCloud_Contact_List_Guide.pptx"
)
prs.save(output_path)
print(f"Saved: {output_path}")
